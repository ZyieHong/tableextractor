from pptx import Presentation
import pandas as pd

def extract_slide_data(ppt_path):
    """提取幻灯片数据（标题+表格）"""
    prs = Presentation(ppt_path)
    slide_info = []
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        # 提取标题
        title = extract_slide_title(slide)
        
        # 提取并处理表格
        tables = process_slide_tables(slide)
        
        slide_info.append({
            "slide_number": slide_num,
            "title": title,
            "tables": tables,
            "table_count": len(tables)
        })
        
        print(f"Slide {slide_num}: {title} [{len(tables)} Table]")
    
    return slide_info

def extract_slide_title(slide):
    """提取幻灯片标题（优化版）"""
    if slide.shapes.title and slide.shapes.title.text.strip():
        return slide.shapes.title.text.strip()
    
    text_shapes = sorted(
        [s for s in slide.shapes if s.has_text_frame],
        key=lambda x: (x.top, x.left)
    )
    
    for shape in text_shapes:
        if shape.text.strip():
            return shape.text.strip()
    
    return "No Title"

def process_slide_tables(slide):
    """处理单张幻灯片中的所有表格"""
    tables = []
    for shape in slide.shapes:
        if shape.has_table:
            table_data = process_table(shape.table)
            df = table_to_dataframe(table_data)
            tables.append(df)
    return tables

def process_table(table):
    """处理单个表格的完整数据"""
    rows = len(table.rows)
    cols = len(table.columns)
    grid = [['' for _ in range(cols)] for _ in range(rows)]
    
    # 处理合并单元格
    for i, row in enumerate(table.rows):
        for j, cell in enumerate(row.cells):
            if cell.is_merge_origin:
                value = cell.text.strip()
                row_span = cell.span_height
                col_span = cell.span_width
                
                # 填充合并区域
                for x in range(i, i + row_span):
                    for y in range(j, j + col_span):
                        if x < rows and y < cols:
                            grid[x][y] = value
            else:
                # 普通单元格直接取值
                grid[i][j] = cell.text.strip()
    
    return grid

def table_to_dataframe(table_data):
    """转换为标准DataFrame"""
    try:
        if not table_data:
            return pd.DataFrame()
        
        # 自动识别表头（首行）
        headers = table_data[0] if table_data else []
        
        # 处理数据行（排除可能的空行）
        data_rows = [row for row in table_data[1:] if any(cell != '' for cell in row)]
        
        # 创建DataFrame
        df = pd.DataFrame(data_rows, columns=headers)
        
        # 清理前导/后缀空格
        df = df.map(lambda x: x.strip() if isinstance(x, str) else x)
        
        return df
    except Exception as e:
        print(f"表格转换异常: {str(e)}")
        return pd.DataFrame(table_data)

# 主程序
if __name__ == "__main__":
    
    keyword = input("Keyword: ").strip()

    # 提取数据
    slides = extract_slide_data(r"C:\Users\User\OneDrive - ump.edu.my\BSD UMP\DSP1111\tests\test_data\cimb.pptx")
    
    # 用户交互
    results = [slide for slide in slides if keyword.lower() in slide["title"].lower()]
    
    if not results:
        print("No Table Found")
    else:
        print(f"\nSlide: {[s['slide_number'] for s in results]}")
        
        for res in results:
            print(f"\n{'='*80}")
            print(f"Slide #{res['slide_number']}")
            print(f"Title: {res['title']}")
            
            if res['table_count'] > 0:
                for idx, df in enumerate(res['tables'], 1):
                    print(f"\n表格 {idx} ({df.shape[1]}列x{df.shape[0]}行)")
                    print(df.to_string(index=False))
                    print("-"*50)
            else:
                print("\nNo Table Found!")
                
                