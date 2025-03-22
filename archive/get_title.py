# 导入依赖
from pptx import Presentation
from pptx.util import Pt
import pandas as pd
from PIL import Image
import pytesseract
import io

# 1. 提取标题和表格状态
from pptx import Presentation
from pptx.util import Pt

def extract_slide_titles(ppt_path):
    prs = Presentation(ppt_path)
    slide_info = []  # Store (slide_num, title, has_table)
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        title = "No Title"
        candidates = []
        
        # Iterate through shapes to filter potential titles
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text.strip()
            if not text:
                continue
            
            # Get font size (take the largest font size)
            max_font_size = 0
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size.pt > max_font_size:
                        max_font_size = run.font.size.pt
            
            # Feature rules: top position + large font + short text
            is_top = shape.top < 914400  # Approximately 1 inch (in EMU units)
            is_large_font = max_font_size > 20
            is_short = len(text) < 50
            
            if is_top and (is_large_font or is_short):
                candidates.append({
                    "text": text,
                    "top": shape.top,
                    "font_size": max_font_size
                })
        
        # Sort by priority: font size > position
        if candidates:
            sorted_candidates = sorted(
                candidates,
                key=lambda x: (-x["font_size"], x["top"])
            )
            title = sorted_candidates[0]["text"]
        else:
            # Method 2: OCR to extract title from background image (fallback)
            title = extract_title_from_background(slide)
        
        # Check if there is a table
        has_table = any(shape.has_table for shape in slide.shapes)
        slide_info.append((slide_num, title, has_table))
        print(f"Slide {slide_num}: Title = {title}, Table = {has_table}")
    
    return slide_info

def extract_title_from_background(slide):
    """Extract title from background image using OCR"""
    if slide.background.fill.type != 1:  # Not an image background
        return "No Title"
    try:
        from PIL import Image
        import pytesseract
        import io
        
        # Extract background image
        image_bytes = slide.background.fill.image.blob
        img = Image.open(io.BytesIO(image_bytes))
        img = img.convert("L")  # Convert to grayscale
        text = pytesseract.image_to_string(img, lang="chi_sim+eng")
        return text.split("\n")[0].strip() if text else "No Title"
    except Exception as e:
        return f"OCR Failed: {str(e)}"

# 2. 根据关键字匹配幻灯片
def find_slides_by_keyword(slide_info, keyword):
    matched = []
    for slide_num, title, has_table in slide_info:
        if keyword.lower() in title.lower():
            matched.append((slide_num, has_table))
    return matched

# Example usage
slide_info = extract_slide_titles("cimb.pptx")
keyword = input("Keyword: ")
matched_slides = find_slides_by_keyword(slide_info, keyword)

if not matched_slides:
    print("No matching slides found.")
else:
    print(f"Matching slides: {[num for num, _ in matched_slides]}")

# 3. 提取表格数据
import pandas as pd

def extract_table_data(slide):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                data.append(row_data)
            return data
    return None

# Extract tables from specified slides
prs = Presentation("cimb.pptx")
for slide_num, has_table in matched_slides:
    slide = prs.slides[slide_num - 1]
    table_data = extract_table_data(slide)
    
    if table_data:
        df = pd.DataFrame(table_data[1:], columns=table_data[0])
        print(f"\nTable data from slide {slide_num}:")
        print(df)
    else:
        print(f"\nNo table found in slide {slide_num}")

# Main program
if __name__ == "__main__":
    ppt_path = "cimb.pptx"
    slide_info = extract_slide_titles(ppt_path)
    
    keyword = input("Keyword: ")
    matched_slides = find_slides_by_keyword(slide_info, keyword)
    
    if not matched_slides:
        print("No matching slides found.")
    else:
        prs = Presentation(ppt_path)
        for slide_num, has_table in matched_slides:
            slide = prs.slides[slide_num - 1]
            table_data = extract_table_data(slide)
            if table_data:
                df = pd.DataFrame(table_data[1:], columns=table_data[0])
                print(f"\nTable data from slide {slide_num}:")
                print(df)
            else:
                print(f"\nNo table found in slide {slide_num}")