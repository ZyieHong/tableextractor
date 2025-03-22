    # from pptx import Presentation
# from pptx.enum.shapes import PP_PLACEHOLDER  # 新增占位符类型检查

# prs = Presentation("cimb.pptx")

# slide_titles = []
# for i, slide in enumerate(prs.slides, start=1):
#     title = "No Title"
    
#     # 新增候选列表（按优先级排序）
#     candidates = []
    
#     # 1. 检查显式标题占位符（最高优先级）
#     if slide.shapes.title and slide.shapes.title.text.strip():
#         candidates.append(slide.shapes.title.text.strip())
    
#     # 2. 检查布局中的标题占位符（新增）
#     if hasattr(slide, 'slide_layout'):
#         for shape in slide.slide_layout.shapes:
#             if shape.is_placeholder and shape.placeholder_format.type in [
#                 PP_PLACEHOLDER.TITLE,
#                 PP_PLACEHOLDER.CENTER_TITLE
#             ]:
#                 candidates.append(shape.text.strip())
    
#     # 3. 检查母版中的标题占位符（新增）
#     if hasattr(slide, 'slide_master'):
#         for shape in slide.slide_master.shapes:
#             if shape.is_placeholder and shape.placeholder_format.type in [
#                 PP_PLACEHOLDER.TITLE,
#                 PP_PLACEHOLDER.CENTER_TITLE
#             ]:
#                 candidates.append(shape.text.strip())
    
#     # 4. 分析其他形状（带防误判机制）
#     for shape in slide.shapes:
#         if shape.has_text_frame and shape.text.strip():
#             text = shape.text.strip()
#             # 防误判条件：
#             if (len(text) < 50 and          # 标题通常较短
#                 any(c.isupper() for c in text) and  # 包含大写字母
#                 shape.width > shape.height * 2):    # 宽高比验证
#                 candidates.append(text)
    
#     # 去重并选择最优候选（保持优先级顺序）
#     if candidates:
#         seen = set()
#         unique_candidates = []
#         for c in candidates:
#             if c not in seen:
#                 seen.add(c)
#                 unique_candidates.append(c)
#         title = unique_candidates[0]

#     slide_titles.append((i, title))
#     print(f"Slide {i} Title: {title}")

# from pptx import Presentation
# from pptx.enum.shapes import PP_PLACEHOLDER

# prs = Presentation("cimb.pptx")

# slide_titles = []
# for i, slide in enumerate(prs.slides, start=1):
#     title = "No Title"
#     candidates = []
    
#     try:
#         # 方法1：显式标题占位符
#         if slide.shapes.title and slide.shapes.title.text.strip():
#             candidates.append(slide.shapes.title.text.strip())
        
#         # 方法2：布局中的标题占位符
#         if hasattr(slide, 'slide_layout'):
#             for shape in slide.slide_layout.shapes:
#                 if shape.is_placeholder and shape.placeholder_format.type in [
#                     PP_PLACEHOLDER.TITLE,
#                     PP_PLACEHOLDER.CENTER_TITLE
#                 ]:
#                     candidates.append(shape.text.strip())
        
#         # 方法3：母版中的标题占位符
#         if hasattr(slide, 'slide_master'):
#             for shape in slide.slide_master.shapes:
#                 if shape.is_placeholder and shape.placeholder_format.type in [
#                     PP_PLACEHOLDER.TITLE,
#                     PP_PLACEHOLDER.CENTER_TITLE
#                 ]:
#                     candidates.append(shape.text.strip())
        
#         # 方法4：形状分析（带验证）
#         for shape in slide.shapes:
#             if shape.has_text_frame and shape.text.strip():
#                 text = shape.text.strip()
#                 try:
#                     if (len(text) < 50 and
#                         any(c.isupper() for c in text) and
#                         shape.width > shape.height * 2 and
#                         not text.lower().startswith("click")):
#                         candidates.append(text)
#                 except AttributeError:
#                     continue
        
#         # 如果所有方法都失败，强制使用第一个形状（不管是否符合验证）
#         if not candidates and slide.shapes:
#             try:
#                 first_shape = slide.shapes[0]
#                 if first_shape.has_text_frame:
#                     candidates.append(first_shape.text.strip())
#             except (AttributeError, IndexError):
#                 pass  # 无形状或无法读取
        
#     except Exception as e:
#         print(f"Slide {i} 处理异常: {str(e)}")
    
#     # 最终决策
#     if candidates:
#         # 去重并取第一个
#         seen = set()
#         unique = [c for c in candidates if c not in seen and not seen.add(c)]
#         title = unique[0]
    
#     slide_titles.append((i, title))
#     print(f"Slide {i} Title: {title}")

# from pptx import Presentation

# prs = Presentation("cimb.pptx")

# slide_titles = []
# for i, slide in enumerate(prs.slides, start=1):
#     title = ""
#     # 优先检查默认标题占位符
#     if slide.shapes.title and slide.shapes.title.text.strip():
#         title = slide.shapes.title.text.strip()
#     else:
#         # 遍历其他形状查找可能的标题
#         for shape in slide.shapes:
#             if shape.has_text_frame and shape.text.strip():
#                 title = shape.text.strip()
#                 break  # 取第一个非空文本框
#     slide_titles.append((i, title))
#     print(f"Slide {i} Title: {title}")

from pptx import Presentation

prs = Presentation("cimb.pptx")

slide_titles = []
for i, slide in enumerate(prs.slides, start=1):
    title = "No Title"
    
    # Priority 1: Check default title placeholder
    if slide.shapes.title and slide.shapes.title.text.strip():
        title = slide.shapes.title.text.strip()
    else:
        # Priority 2: Check other shapes sorted by vertical position (top first)
        text_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
        sorted_shapes = sorted(text_shapes, key=lambda x: x.top)  # Sort by top position
        
        for shape in sorted_shapes:
            text = shape.text.strip()
            if text:
                title = text
                break  # Take the topmost non-empty text box
    
    slide_titles.append((i, title))
    print(f"Slide {i} Title: {title}")
    # print(f'[{i}], ["{title}"]')