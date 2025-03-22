from pptx import Presentation
import pandas as pd

def extract_slide_data(ppt_path):
    """提取幻灯片数据（标题+表格）"""
    prs = Presentation(ppt_path)
    slide_info = []
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        title = _extract_slide_title(slide)
        tables = _process_slide_tables(slide)
        
        slide_info.append({
            "slide_number": slide_num,
            "title": title,
            "tables": tables,
            "table_count": len(tables)
        })
    
    return slide_info

def _extract_slide_title(slide):
    """提取幻灯片标题（内部方法）"""
    # 优先使用官方标题占位符
    if slide.shapes.title and slide.shapes.title.text.strip():
        return slide.shapes.title.text.strip()
    
    # 备用方案：按位置排序的文本形状
    text_shapes = sorted(
        [s for s in slide.shapes if s.has_text_frame],
        key=lambda x: (x.top, x.left)
    )
    
    for shape in text_shapes:
        if shape.text.strip():
            return shape.text.strip()
    
    return "No Title"

def _process_slide_tables(slide):
    """处理单张幻灯片中的所有表格（内部方法）"""
    tables = []
    for shape in slide.shapes:
        if shape.has_table:
            table_data = _process_table(shape.table)
            df = _table_to_dataframe(table_data)
            tables.append(df)
    return tables

def _process_table(table):
    """处理单个表格的完整数据（内部方法）"""
    rows = len(table.rows)
    cols = len(table.columns)
    grid = [['' for _ in range(cols)] for _ in range(rows)]
    
    # 处理合并单元格
    for i, row in enumerate(table.rows):
        for j, cell in enumerate(row.cells):
            if cell.is_merge_origin:
                value = cell.text.strip()
                row_span = cell.span_height
                col_span = cell.span_width
                
                # 填充合并区域
                for x in range(i, i + row_span):
                    for y in range(j, j + col_span):
                        if x < rows and y < cols:
                            grid[x][y] = value
            else:
                # 普通单元格直接取值（避免覆盖合并单元格）
                if not grid[i][j]:
                    grid[i][j] = cell.text.strip()
    
    return grid

def _table_to_dataframe(table_data):
    """转换为标准DataFrame（内部方法）"""
    try:
        if not table_data or len(table_data) < 1:
            return pd.DataFrame()
        
        # 自动识别表头（首行非空行）
        header_row = next(
            (i for i, row in enumerate(table_data) if any(cell.strip() for cell in row)),
            None
        )
        
        if header_row is None:
            return pd.DataFrame()
        
        headers = [cell.strip() for cell in table_data[header_row]]
        
        # 处理数据行（排除可能的空行）
        data_rows = [
            [cell.strip() for cell in row]
            for row in table_data[header_row+1:]
            if any(cell.strip() for cell in row)
        ]
        
        # 创建DataFrame
        df = pd.DataFrame(data_rows, columns=headers)
        
        # 清理空白值
        df.replace('', pd.NA, inplace=True)
        df.dropna(how='all', inplace=True)
        
        return df
    
    except Exception as e:
        print(f"表格转换异常: {str(e)}")
        return pd.DataFrame()